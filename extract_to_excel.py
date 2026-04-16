"""
extract_to_excel.py
-------------------
Extracts raw text from every page/slide of PDFs and PPTX files,
generates a one-sentence summary per page via the Anthropic API,
and writes everything to a structured Excel file.

Columns: Source File | Page Number | Raw Text | Summary

Install dependencies:
    pip install pdfplumber python-pptx openpyxl anthropic

Usage:
    python extract_to_excel.py file1.pdf deck.pptx ...
    python extract_to_excel.py                          # auto-detects all PDFs/PPTXs in current folder
    python extract_to_excel.py file.pdf --pages 20-25  # only process pages 20 to 25
"""

import sys
import os
import re
import glob
import argparse
import anthropic
import pdfplumber
from pptx import Presentation
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.utils import get_column_letter


# ── Config ────────────────────────────────────────────────────────────────────

OUTPUT_FILE = "extracted_insights.xlsx"
MODEL       = "claude-sonnet-4-6"

SUMMARY_PROMPT = (
    "Summarize the following page/slide content in 1-2 concise sentences. "
    "Focus on the key insight or main point.\n\nContent:\n{text}"
)


# ── Text sanitization ─────────────────────────────────────────────────────────

# openpyxl rejects control characters (except tab/newline/CR)
_ILLEGAL_CHARS_RE = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f]")

def sanitize(text: str) -> str:
    text = _ILLEGAL_CHARS_RE.sub("", text)
    text = text.replace("\u26ab", "\u2022")  # ⚫ → •
    return text


# ── Extraction helpers ────────────────────────────────────────────────────────

def parse_page_selection(value: str) -> set[int]:
    """
    Parse a flexible page selection string into a set of page numbers.

    Supports:
        "5"          → {5}
        "1-5"        → {1, 2, 3, 4, 5}
        "1,3,5"      → {1, 3, 5}
        "1-5,10,15-18" → {1,2,3,4,5,10,15,16,17,18}
    """
    selected = set()
    for part in value.split(","):
        part = part.strip()
        if "-" in part:
            try:
                start, end = part.split("-", 1)
                selected.update(range(int(start), int(end) + 1))
            except ValueError:
                raise argparse.ArgumentTypeError(
                    f"Invalid range '{part}'. Use format START-END (e.g. 20-25)."
                )
        else:
            try:
                selected.add(int(part))
            except ValueError:
                raise argparse.ArgumentTypeError(
                    f"Invalid page number '{part}'. Must be an integer."
                )
    return selected


def extract_pdf_pages(path: str, page_selection: set[int] | None = None) -> list[dict]:
    """Return list of {page_num, raw_text} for PDF pages, optionally filtered by selection."""
    pages = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            if page_selection and i not in page_selection:
                continue
            text = page.extract_text() or ""
            pages.append({"page_num": i, "raw_text": text.strip()})
    return pages


def extract_pptx_slides(path: str, page_selection: set[int] | None = None) -> list[dict]:
    """Return list of {page_num, raw_text} for PPTX slides, optionally filtered by selection."""
    slides = []
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, start=1):
        if page_selection and i not in page_selection:
            continue
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        slides.append({"page_num": i, "raw_text": "\n".join(texts)})
    return slides


def get_summary(client: anthropic.Anthropic, text: str) -> str:
    """Call Claude to summarize a single page/slide."""
    if not text.strip():
        return "(no text on this page)"
    msg = client.messages.create(
        model=MODEL,
        max_tokens=200,
        messages=[{"role": "user", "content": SUMMARY_PROMPT.format(text=text[:3000])}],
    )
    return msg.content[0].text.strip()


# ── Excel writer ──────────────────────────────────────────────────────────────

HEADER_FILL   = PatternFill("solid", start_color="1F4E79")   # dark blue
HEADER_FONT   = Font(name="Arial", bold=True, color="FFFFFF", size=11)
ALT_FILL      = PatternFill("solid", start_color="D6E4F0")   # light blue
NORMAL_FILL   = PatternFill("solid", start_color="FFFFFF")
CELL_FONT     = Font(name="Arial", size=10)
WRAP          = Alignment(wrap_text=True, vertical="top")
CENTER_TOP    = Alignment(horizontal="center", vertical="top")

COL_WIDTHS = {
    "A": 30,   # Source File
    "B": 14,   # Page Number
    "C": 70,   # Raw Text
    "D": 55,   # Summary
}

HEADERS = ["Source File", "Page Number", "Raw Text", "Summary"]


def write_excel(rows: list[dict], output_path: str):
    wb = Workbook()
    ws = wb.active
    ws.title = "Extracted Insights"

    # Header row
    for col, header in enumerate(HEADERS, start=1):
        cell = ws.cell(row=1, column=col, value=header)
        cell.font   = HEADER_FONT
        cell.fill   = HEADER_FILL
        cell.alignment = Alignment(horizontal="center", vertical="center")

    ws.row_dimensions[1].height = 22

    # Data rows
    for row_idx, row in enumerate(rows, start=2):
        fill = ALT_FILL if row_idx % 2 == 0 else NORMAL_FILL

        def write(col, value, align=WRAP):
            c = ws.cell(row=row_idx, column=col, value=value)
            c.font      = CELL_FONT
            c.fill      = fill
            c.alignment = align

        write(1, row["source_file"])
        write(2, row["page_num"], CENTER_TOP)
        write(3, sanitize(row["raw_text"]))
        write(4, sanitize(row["summary"]))

    # Column widths
    for letter, width in COL_WIDTHS.items():
        ws.column_dimensions[letter].width = width

    # Freeze header row
    ws.freeze_panes = "A2"

    wb.save(output_path)
    print(f"\n✅  Saved → {output_path}")


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Extract PDF/PPTX pages to Excel with AI summaries.")
    parser.add_argument("files", nargs="*", help="PDF or PPTX files to process")
    parser.add_argument(
        "--pages",
        type=parse_page_selection,
        metavar="SELECTION",
        help=(
            "Pages to process. Supports individual pages, ranges, or a mix. "
            "Examples: --pages 5  |  --pages 1-10  |  --pages 1-5,10,15-18"
        ),
    )
    args = parser.parse_args()

    # Collect input files
    files = args.files if args.files else glob.glob("*.pdf") + glob.glob("*.pptx")
    page_selection = args.pages

    if not files:
        print("No PDF or PPTX files found. Pass file paths as arguments or run from the folder containing them.")
        sys.exit(1)

    if page_selection:
        display = ", ".join(str(p) for p in sorted(page_selection))
        print(f"   Page selection: {display}")

    client = anthropic.Anthropic()   # reads ANTHROPIC_API_KEY from environment
    all_rows = []

    for file_path in files:
        ext = os.path.splitext(file_path)[1].lower()
        name = os.path.basename(file_path)

        if ext == ".pdf":
            pages = extract_pdf_pages(file_path, page_selection)
        elif ext in (".pptx", ".ppt"):
            pages = extract_pptx_slides(file_path, page_selection)
        else:
            print(f"  ⚠  Skipping unsupported file: {file_path}")
            continue

        print(f"\n📄  Processing {name}  ({len(pages)} pages/slides)")

        for page in pages:
            print(f"   Page {page['page_num']:>3} — summarising…", end="\r")
            summary = get_summary(client, page["raw_text"])
            all_rows.append({
                "source_file": name,
                "page_num":    page["page_num"],
                "raw_text":    page["raw_text"],
                "summary":     summary,
            })

        print(f"   ✓  {len(pages)} pages done{' ' * 20}")

    if not all_rows:
        print("No content extracted.")
        sys.exit(1)

    write_excel(all_rows, OUTPUT_FILE)
    print(f"   Total rows written: {len(all_rows)}")


if __name__ == "__main__":
    main()
