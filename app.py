"""
app.py
------
Streamlit UI for extract_to_excel.py.
Upload a PDF or PPTX, optionally select pages, and download the Excel summary.
"""

import io
import os
import re
import tempfile

import anthropic
import pdfplumber
import streamlit as st
from dotenv import load_dotenv
from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill
from pptx import Presentation

load_dotenv()

# ── Config ────────────────────────────────────────────────────────────────────

MODEL = "claude-sonnet-4-6"

SUMMARY_PROMPT = (
    "Summarize the following page/slide content in 1-2 concise sentences. "
    "Focus on the key insight or main point.\n\nContent:\n{text}"
)

# ── Text sanitization ─────────────────────────────────────────────────────────

_ILLEGAL_CHARS_RE = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f]")


def sanitize(text: str) -> str:
    text = _ILLEGAL_CHARS_RE.sub("", text)
    text = text.replace("\u26ab", "\u2022")  # ⚫ → •
    return text


# ── Page selection parser ─────────────────────────────────────────────────────

def parse_page_selection(value: str) -> set[int] | None:
    """
    Parse a flexible page selection string into a set of page numbers.
    Returns None if the string is empty (meaning: all pages).
    Supports: "5", "1-5", "1,3,5", "1-5,10,15-18"
    """
    value = value.strip()
    if not value:
        return None
    selected = set()
    for part in value.split(","):
        part = part.strip()
        if "-" in part:
            start, end = part.split("-", 1)
            selected.update(range(int(start), int(end) + 1))
        else:
            selected.add(int(part))
    return selected


# ── Extraction helpers ────────────────────────────────────────────────────────

def extract_pdf_pages(path: str, page_selection: set[int] | None) -> list[dict]:
    pages = []
    with pdfplumber.open(path) as pdf:
        total = len(pdf.pages)
        for i, page in enumerate(pdf.pages, start=1):
            if page_selection and i not in page_selection:
                continue
            text = page.extract_text() or ""
            pages.append({"page_num": i, "total": total, "raw_text": text.strip()})
    return pages


def extract_pptx_slides(path: str, page_selection: set[int] | None) -> list[dict]:
    slides = []
    prs = Presentation(path)
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, start=1):
        if page_selection and i not in page_selection:
            continue
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        slides.append({"page_num": i, "total": total, "raw_text": "\n".join(texts)})
    return slides


def get_summary(client: anthropic.Anthropic, text: str) -> str:
    if not text.strip():
        return "(no text on this page)"
    msg = client.messages.create(
        model=MODEL,
        max_tokens=200,
        messages=[{"role": "user", "content": SUMMARY_PROMPT.format(text=text[:3000])}],
    )
    return msg.content[0].text.strip()


# ── Excel builder ─────────────────────────────────────────────────────────────

HEADER_FILL = PatternFill("solid", start_color="1F4E79")
HEADER_FONT = Font(name="Arial", bold=True, color="FFFFFF", size=11)
ALT_FILL    = PatternFill("solid", start_color="D6E4F0")
NORMAL_FILL = PatternFill("solid", start_color="FFFFFF")
CELL_FONT   = Font(name="Arial", size=10)
WRAP        = Alignment(wrap_text=True, vertical="top")
CENTER_TOP  = Alignment(horizontal="center", vertical="top")

COL_WIDTHS = {"A": 30, "B": 14, "C": 70, "D": 55}
HEADERS    = ["Source File", "Page Number", "Raw Text", "Summary"]


def build_excel(rows: list[dict]) -> bytes:
    wb = Workbook()
    ws = wb.active
    ws.title = "Extracted Insights"

    for col, header in enumerate(HEADERS, start=1):
        cell = ws.cell(row=1, column=col, value=header)
        cell.font      = HEADER_FONT
        cell.fill      = HEADER_FILL
        cell.alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 22

    for row_idx, row in enumerate(rows, start=2):
        fill = ALT_FILL if row_idx % 2 == 0 else NORMAL_FILL

        def write(col, value, align=WRAP):
            c = ws.cell(row=row_idx, column=col, value=value)
            c.font = CELL_FONT
            c.fill = fill
            c.alignment = align

        write(1, row["source_file"])
        write(2, row["page_num"], CENTER_TOP)
        write(3, sanitize(row["raw_text"]))
        write(4, sanitize(row["summary"]))

    for letter, width in COL_WIDTHS.items():
        ws.column_dimensions[letter].width = width
    ws.freeze_panes = "A2"

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ── Streamlit UI ──────────────────────────────────────────────────────────────

st.set_page_config(page_title="Distill Insights", page_icon="📄", layout="centered")

st.title("📄 Distill Insights")
st.caption("Upload a PDF or PPTX — get an Excel file with raw text and AI summaries per page.")

# API key check
api_key = os.getenv("ANTHROPIC_API_KEY")
if not api_key:
    api_key = st.text_input("Anthropic API Key", type="password", placeholder="sk-ant-...")
    if not api_key:
        st.stop()

# File upload
uploaded = st.file_uploader("Upload file", type=["pdf", "pptx", "ppt"])

if uploaded:
    ext = os.path.splitext(uploaded.name)[1].lower()

    # Write to temp file so pdfplumber / pptx can read it
    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        tmp.write(uploaded.read())
        tmp_path = tmp.name

    # Get total page count for context
    try:
        if ext == ".pdf":
            with pdfplumber.open(tmp_path) as pdf:
                total_pages = len(pdf.pages)
        else:
            prs = Presentation(tmp_path)
            total_pages = len(prs.slides)
        st.info(f"**{uploaded.name}** — {total_pages} {'pages' if ext == '.pdf' else 'slides'} total")
    except Exception as e:
        st.error(f"Could not read file: {e}")
        st.stop()

    # Page selection
    page_input = st.text_input(
        "Page selection (leave blank for all pages)",
        placeholder="e.g.  5   or   1-10   or   1-5, 10, 15-18",
    )

    # Validate input before running
    page_selection = None
    input_error = None
    if page_input.strip():
        try:
            page_selection = parse_page_selection(page_input)
            out_of_range = {p for p in page_selection if p < 1 or p > total_pages}
            if out_of_range:
                input_error = f"Pages out of range: {sorted(out_of_range)} — file only has {total_pages} pages."
            else:
                st.caption(f"Selected: pages {', '.join(str(p) for p in sorted(page_selection))}")
        except ValueError:
            input_error = "Invalid format. Use numbers, ranges (1-5), or a mix (1-5, 10, 15-18)."

    if input_error:
        st.error(input_error)
        st.stop()

    # Run button
    if st.button("Extract & Summarize", type="primary"):
        client = anthropic.Anthropic(api_key=api_key)

        if ext == ".pdf":
            pages = extract_pdf_pages(tmp_path, page_selection)
        else:
            pages = extract_pptx_slides(tmp_path, page_selection)

        if not pages:
            st.warning("No pages matched the selection.")
            st.stop()

        all_rows = []
        progress = st.progress(0, text="Starting…")

        for idx, page in enumerate(pages):
            progress.progress(idx / len(pages), text=f"Summarizing page {page['page_num']}…")
            summary = get_summary(client, page["raw_text"])
            all_rows.append({
                "source_file": uploaded.name,
                "page_num":    page["page_num"],
                "raw_text":    page["raw_text"],
                "summary":     summary,
            })

        progress.progress(1.0, text="Done!")

        excel_bytes = build_excel(all_rows)
        base_name = os.path.splitext(uploaded.name)[0]

        st.success(f"Extracted {len(all_rows)} pages. Download below.")
        st.download_button(
            label="Download Excel",
            data=excel_bytes,
            file_name=f"{base_name}_insights.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )

    os.unlink(tmp_path)
